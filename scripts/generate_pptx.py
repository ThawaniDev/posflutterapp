import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def parse_markdown(filepath):
    with open(filepath, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Split content by '---'
    sections = content.split('---')
    slides = []
    
    for sec in sections:
        sec = sec.strip()
        if not sec:
            continue
            
        lines = sec.split('\n')
        # Check if any heading in this section contains 'الشريحة'
        has_slide_heading = False
        for line in lines:
            if line.strip().startswith('##') and 'الشريحة' in line:
                has_slide_heading = True
                break
        if not has_slide_heading:
            continue
        slide_info = {
            'slide_num_raw': '',
            'slide_label': '',
            'title': '',
            'subtitle': '',
            'bullets': [],
            'visual': '',
            'presenter_note': '',
            'guide': ''
        }
        
        # Parse the heading: e.g. ## الشريحة 1 - الغلاف
        for line in lines:
            if line.strip().startswith('## '):
                heading = line.replace('##', '').strip()
                slide_info['slide_label'] = heading
                # Extract number
                num_match = re.search(r'\d+', heading)
                if num_match:
                    slide_info['slide_num_raw'] = num_match.group()
                break
                
        # Parse fields line by line
        current_field = None
        for line in lines:
            line_str = line.strip()
            if not line_str:
                continue
            if line_str.startswith('## '):
                continue
                
            if line_str.startswith('العنوان:'):
                slide_info['title'] = line_str.replace('العنوان:', '', 1).strip()
                current_field = 'title'
            elif line_str.startswith('العنوان الفرعي:'):
                slide_info['subtitle'] = line_str.replace('العنوان الفرعي:', '', 1).strip()
                current_field = 'subtitle'
            elif line_str.startswith('النص الرئيسي:'):
                current_field = 'bullets'
            elif line_str.startswith('الاتجاه البصري:'):
                slide_info['visual'] = line_str.replace('الاتجاه البصري:', '', 1).strip()
                current_field = 'visual'
            elif line_str.startswith('ملاحظة العرض:'):
                slide_info['presenter_note'] = line_str.replace('ملاحظة العرض:', '', 1).strip()
                current_field = 'presenter_note'
            elif line_str.startswith('الدليل:'):
                slide_info['guide'] = line_str.replace('الدليل:', '', 1).strip()
                current_field = 'guide'
            else:
                # Handle multiline flows
                if current_field == 'bullets' and line_str.startswith('-'):
                    bullet_text = line_str.lstrip('- ').strip()
                    slide_info['bullets'].append(bullet_text)
                elif current_field == 'bullets' and not line_str.startswith('-'):
                    prefixes = ['العنوان:', 'العنوان الفرعي:', 'الاتجاه البصري:', 'ملاحظة العرض:', 'الدليل:', '## ']
                    if not any(line_str.startswith(p) for prefix in prefixes for p in [prefix]):
                        if slide_info['bullets']:
                            slide_info['bullets'][-1] += " " + line_str
                elif current_field == 'subtitle':
                    prefixes = ['العنوان:', 'النص الرئيسي:', 'الاتجاه البصري:', 'ملاحظة العرض:', 'الدليل:', '## ']
                    if not any(line_str.startswith(prefix) for prefix in prefixes):
                        slide_info['subtitle'] += " " + line_str
                elif current_field == 'visual':
                    prefixes = ['العنوان:', 'النص الرئيسي:', 'العنوان الفرعي:', 'ملاحظة العرض:', 'الدليل:', '## ']
                    if not any(line_str.startswith(prefix) for prefix in prefixes):
                        slide_info['visual'] += " " + line_str
                elif current_field == 'presenter_note':
                    prefixes = ['العنوان:', 'النص الرئيسي:', 'العنوان الفرعي:', 'الاتجاه البصري:', 'الدليل:', '## ']
                    if not any(line_str.startswith(prefix) for prefix in prefixes):
                        slide_info['presenter_note'] += " " + line_str
                elif current_field == 'guide':
                    prefixes = ['العنوان:', 'النص الرئيسي:', 'العنوان الفرعي:', 'الاتجاه البصري:', 'ملاحظة العرض:', '## ']
                    if not any(line_str.startswith(prefix) for prefix in prefixes):
                        slide_info['guide'] += " " + line_str
                        
        slides.append(slide_info)
    return slides

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_shape(slide, shape_type, left, top, width, height, fill_color, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_paragraph_formatted(text_frame, text, font_size=14, color=RGBColor(15, 23, 42), bold=False, italic=False, align=PP_ALIGN.RIGHT):
    p = text_frame.add_paragraph() if len(text_frame.text) > 0 else text_frame.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = "Cairo"
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    return p

def add_header(slide, slide_no_str, title_text, subtitle_text, dark_mode=False):
    title_box = slide.shapes.add_textbox(Inches(0.66), Inches(0.4), Inches(12.0), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    # Slide Number prefix indicator (e.g. "05 / 20")
    num_color = RGBColor(255, 191, 13) if dark_mode else RGBColor(253, 130, 9)
    p_num = tf.paragraphs[0]
    p_num.text = f"{slide_no_str.zfill(2)}  |"
    p_num.alignment = PP_ALIGN.RIGHT
    p_num.font.name = "Cairo"
    p_num.font.size = Pt(12)
    p_num.font.bold = True
    p_num.font.color.rgb = num_color
    
    # Title
    t_color = RGBColor(255, 255, 255) if dark_mode else RGBColor(15, 23, 42)
    p_title = tf.add_paragraph()
    p_title.text = title_text
    p_title.alignment = PP_ALIGN.RIGHT
    p_title.font.name = "Cairo"
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = t_color
    
    # Subtitle
    if subtitle_text:
        s_color = RGBColor(248, 247, 245) if dark_mode else RGBColor(71, 85, 105)
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle_text
        p_sub.alignment = PP_ALIGN.RIGHT
        p_sub.font.name = "Cairo"
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = s_color

def set_speaker_notes(slide, slide_data):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    notes_text = []
    
    if slide_data.get('visual'):
        notes_text.append(f"Visual Cue (الاتجاه البصري):\n{slide_data['visual']}\n")
    if slide_data.get('presenter_note'):
        notes_text.append(f"Presenter Notes (ملاحظة العرض):\n{slide_data['presenter_note']}\n")
    if slide_data.get('guide'):
        notes_text.append(f"Code Reference Guide (الدليل المرجعي):\n{slide_data['guide']}")
        
    text_frame.text = "\n".join(notes_text)

def main():
    # Palette definitions
    ORANGE = RGBColor(253, 130, 9)
    YELLOW = RGBColor(255, 191, 13)
    DARK_BLUE = RGBColor(15, 23, 42)
    MUTED_GREY = RGBColor(71, 85, 105)
    BG_LIGHT = RGBColor(248, 247, 245)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_ORANGE = RGBColor(254, 232, 217)
    
    # Resolve project paths
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    markdown_path = os.path.join(base_dir, 'provider', 'pos_provider_brochure_slide_deck_ar.md')
    output_path = os.path.join(base_dir, 'provider', 'pos_provider_brochure_slide_deck_ar.pptx')
    
    print(f"Reading markdown from: {markdown_path}")
    slides_data = parse_markdown(markdown_path)
    print(f"Parsed {len(slides_data)} slides from markdown.")
    
    prs = Presentation()
    # Force 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    for i, sd in enumerate(slides_data):
        slide_no = sd['slide_num_raw'] if sd['slide_num_raw'] else str(i + 1)
        slide_title = sd['title']
        slide_sub = sd['subtitle']
        bullets = sd['bullets']
        
        slide = prs.slides.add_slide(blank_layout)
        
        # Default Background
        set_slide_background(slide, BG_LIGHT)
        
        # Set speaker notes
        set_speaker_notes(slide, sd)
        
        # 1. Slide 1 (Cover Slide)
        if slide_no == '1':
            set_slide_background(slide, DARK_BLUE)
            
            # Decorative abstract orange/yellow shapes
            create_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5), ORANGE)
            create_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0), Inches(0.15), Inches(7.5), YELLOW)
            
            # Massive cover textbox (Centered/Right-heavy)
            cover_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.5), Inches(4.5))
            tf = cover_box.text_frame
            tf.word_wrap = True
            
            # Tag/Category
            add_paragraph_formatted(tf, "عائلة منتجات WAMEED & THAWANI", 14, YELLOW, bold=True)
            tf.add_paragraph() # spacing
            
            # Title
            add_paragraph_formatted(tf, slide_title, 36, WHITE, bold=True)
            tf.add_paragraph()
            
            # Subtitle
            add_paragraph_formatted(tf, slide_sub, 16, BG_LIGHT, bold=False)
            tf.add_paragraph()
            
            # Small tagline / details
            tagline = "بروشور تشغيلي متكامل للمزوّدين ومحطات البيع وإدارة القنوات الإلكترونية"
            add_paragraph_formatted(tf, tagline, 12, ORANGE, bold=False, italic=True)
            
        # 2. Slide 20 (Final Closing Slide)
        elif slide_no == '20':
            set_slide_background(slide, DARK_BLUE)
            
            # Colored details border
            create_shape(slide, MSO_SHAPE.RECTANGLE, Inches(12.933), Inches(0), Inches(0.4), Inches(7.5), ORANGE)
            create_shape(slide, MSO_SHAPE.RECTANGLE, Inches(12.783), Inches(0), Inches(0.15), Inches(7.5), YELLOW)
            
            cover_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(10.5), Inches(4.5))
            tf = cover_box.text_frame
            tf.word_wrap = True
            
            # Closing Tag
            add_paragraph_formatted(tf, "مجموعة تشغيل المتجر وإدارة النمو المتصل", 14, YELLOW, bold=True)
            tf.add_paragraph()
            
            # Title
            add_paragraph_formatted(tf, slide_title, 34, WHITE, bold=True)
            tf.add_paragraph()
            
            # Subtitle
            add_paragraph_formatted(tf, slide_sub, 15, BG_LIGHT, bold=False)
            tf.add_paragraph()
            
            # Small footnote
            footnote = "تم التطوير والدعم الفني والربط المباشر ضمن منصة موحدة مبنية على Flutter"
            add_paragraph_formatted(tf, footnote, 12, ORANGE, bold=False, italic=True)
            
        # 3. Slide 2 (Clean 4-column corporate overview)
        elif slide_no == '2':
            add_header(slide, slide_no, slide_title, slide_sub)
            
            # Slide 2 expects 4 columns for: البيع، التشغيل، الامتثال، الاتصال
            # There are 4 bullets: Bullet 1 is Col 4 (right), Bullet 2 is Col 3, Bullet 3 is Col 2, Bullet 4 is Col 1 (left)
            col_width = Inches(2.8)
            col_gap = Inches(0.2)
            top_y = Inches(1.9)
            height_y = Inches(4.8)
            
            col_names = ["البيع والـ POS (1)", "العمق التشغيلي (2)", "الطلب الإلكتروني (3)", "الأساس الموثوق (4)"]
            # Swap order for RTL reading: Box 1 is at index 3 (Rightmost), Box 4 is at index 0 (Leftmost)
            column_positions = [
                Inches(0.66),                            # Left-most col (Col 4)
                Inches(0.66 + 3.0),                      # Second col (Col 3)
                Inches(0.66 + 6.0),                      # Third col (Col 2)
                Inches(0.66 + 9.0)                       # Right-most col (Col 1)
            ]
            
            # Reverse bullets to match RTL Columns
            r_bullets = list(reversed(bullets)) if len(bullets) >= 4 else bullets + [""] * (4 - len(bullets))
            r_names = list(reversed(col_names))
            
            for j in range(4):
                left_x = column_positions[j]
                
                # Card shape
                create_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left_x, top_y, col_width, height_y, WHITE, MUTED_GREY, 1)
                
                # Visual Indicator Strip
                header_strip_color = ORANGE if j in [2, 3] else DARK_BLUE # Highlight POS on direct right
                create_shape(slide, MSO_SHAPE.RECTANGLE, left_x + Inches(0.1), top_y + Inches(0.1), col_width - Inches(0.2), Inches(0.15), header_strip_color)
                
                # Card Text
                col_box = slide.shapes.add_textbox(left_x + Inches(0.1) , top_y + Inches(0.4), col_width - Inches(0.2), height_y - Inches(0.5))
                ctf = col_box.text_frame
                ctf.word_wrap = True
                
                col_title = r_names[j]
                add_paragraph_formatted(ctf, col_title, 14, header_strip_color, bold=True)
                ctf.add_paragraph() # space
                
                bullet_text = r_bullets[j]
                add_paragraph_formatted(ctf, bullet_text, 11, DARK_BLUE, bold=False)

        # 4. Slide 3 (2-column layout + small common services bottom section)
        elif slide_no == '3':
            add_header(slide, slide_no, slide_title, slide_sub)
            
            # Left block: Wameed POS. Right block: Thawani Connected Commerce. Bottom row: Common Services
            # Bullets: Bullet 1 is Thawani, Bullet 2 is Wameed, Bullet 3 is Common Services, Bullet 4 is guide
            col_width = Inches(5.8)
            top_y = Inches(1.8)
            height_y = Inches(3.6)
            
            # Thawani (Right Column in RTL) - Col 1 is on the right, Col 2 on the left
            left_pos_thawani = Inches(6.86)
            create_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left_pos_thawani, top_y, col_width, height_y, WHITE, ORANGE, 1)
            create_shape(slide, MSO_SHAPE.RECTANGLE, left_pos_thawani + Inches(0.15), top_y + Inches(0.15), col_width - Inches(0.3), Inches(0.2), ORANGE)
            
            thawani_box = slide.shapes.add_textbox(left_pos_thawani + Inches(0.15), top_y + Inches(0.45), col_width - Inches(0.3), height_y - Inches(0.6))
            ttf = thawani_box.text_frame
            ttf.word_wrap = True
            add_paragraph_formatted(ttf, "المنتج الأول: Thawani Connected Commerce", 15, ORANGE, bold=True)
            text_t = bullets[0] if len(bullets) > 0 else ""
            add_paragraph_formatted(ttf, text_t, 11, DARK_BLUE)
            
            # Wameed POS (Left Column in RTL)
            left_pos_wameed = Inches(0.66)
            create_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left_pos_wameed, top_y, col_width, height_y, WHITE, DARK_BLUE, 1)
            create_shape(slide, MSO_SHAPE.RECTANGLE, left_pos_wameed + Inches(0.15), top_y + Inches(0.15), col_width - Inches(0.3), Inches(0.2), DARK_BLUE)
            
            wameed_box = slide.shapes.add_textbox(left_pos_wameed + Inches(0.15), top_y + Inches(0.45), col_width - Inches(0.3), height_y - Inches(0.6))
            wtf = wameed_box.text_frame
            wtf.word_wrap = True
            add_paragraph_formatted(wtf, "المنتج الثاني: Wameed POS", 15, DARK_BLUE, bold=True)
            text_w = bullets[1] if len(bullets) > 1 else ""
            add_paragraph_formatted(wtf, text_w, 11, DARK_BLUE)
            
            # Bottom section (Common services across the screen)
            bottom_y = Inches(5.6)
            bottom_height = Inches(1.3)
            create_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.66), bottom_y, Inches(12.0), bottom_height, L_BG := RGBColor(240, 240, 240), DARK_BLUE, 1)
            
            bottom_box = slide.shapes.add_textbox(Inches(0.86), bottom_y + Inches(0.15), Inches(11.6), bottom_height - Inches(0.3))
            btf = bottom_box.text_frame
            btf.word_wrap = True
            add_paragraph_formatted(btf, "الأساس والخدمات الموحدة المشتركة", 13, DARK_BLUE, bold=True)
            text_common = bullets[2] if len(bullets) > 2 else "اختيار الفرع، التراخيص، الصلاحيات، الإشعارات، التوطين المشترك، ومستويات الوصول الدعم الفني."
            add_paragraph_formatted(btf, text_common, 11, MUTED_GREY)

        # 5. Slide 4 (Horizontal timeline / Operating model)
        elif slide_no == '4':
            add_header(slide, slide_no, slide_title, slide_sub)
            
            # 3 horizontal boxes representing timeline of operations, right to left.
            box_width = Inches(3.7)
            top_y = Inches(2.2)
            height_y = Inches(4.3)
            col_gap = Inches(0.4)
            
            # RTL sequence: Step 1 is rightmost, Step 2 is middle, Step 3 is leftmost
            positions = [
                Inches(0.66 + 2 * (3.7 + 0.4)),  # Rightmost: Step 1 (Store floor POS)
                Inches(0.66 + 3.7 + 0.4),        # Middle: Step 2 (Staff & Management)
                Inches(0.66)                     # Leftmost: Step 3 (Connected Commerce)
            ]
            
            # Gather timeline texts from slide 4 bullets
            t_bullets = bullets + [""] * (3 - len(bullets))
            step_titles = [
                "1. أرضية المتجر والبيع المباشر",
                "2. التحكم المتكامل والخصائص المشتركة",
                "3. التجارة الإلكترونية المتصلة"
            ]
            
            colors = [DARK_BLUE, MUTED_GREY, ORANGE]
            
            for j in range(3):
                left_x = positions[j]
                
                # Card
                create_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left_x, top_y, box_width, height_y, WHITE, colors[j], 1)
                # Accent Header Stripe
                create_shape(slide, MSO_SHAPE.RECTANGLE, left_x + Inches(0.1), top_y + Inches(0.1), box_width - Inches(0.2), Inches(0.15), colors[j])
                
                box_txt = slide.shapes.add_textbox(left_x + Inches(0.1), top_y + Inches(0.35), box_width - Inches(0.2), height_y - Inches(0.45))
                tf_box = box_txt.text_frame
                tf_box.word_wrap = True
                
                add_paragraph_formatted(tf_box, step_titles[j], 14, colors[j], bold=True)
                tf_box.add_paragraph() # space
                
                add_paragraph_formatted(tf_box, t_bullets[j], 11, DARK_BLUE)
                
                # Draw connecting chevron visual for steps (RTL: draw arrow pointing to the left!)
                if j < 2: # No arrow on left-most step
                    arrow_x = left_x - Inches(0.3)
                    arrow_y = top_y + Inches(2.0)
                    create_shape(slide, MSO_SHAPE.LEFT_ARROW, arrow_x, arrow_y, Inches(0.2), Inches(0.3), YELLOW)

        # 6. Slide 7 (Thawani Market share KPI cards)
        elif slide_no == '7':
            add_header(slide, slide_no, slide_title, slide_sub)
            
            # Draw KPI Cards (3 boxes) showing target metrics
            box_width = Inches(3.7)
            top_y = Inches(1.9)
            height_y = Inches(4.8)
            
            # RTL Layout: Rightmost = Entry, Middle = 5% Target, Leftmost = Midterm Target
            positions = [
                Inches(8.86),  # Right
                Inches(4.76),  # Middle
                Inches(0.66)   # Left
            ]
            
            kpi_metrics = ["اطلاق سريع", "السنة الثانية ~5%", "المدى المتوسط 5-8%"]
            kpi_titles = ["الفرصة وسرعة الدخول", "الحصة المستهدفة من التوصيل", "الأفق والنمو التشغيلي"]
            kpi_texts = [
                "تأسيس القنوات الإلكترونية وربط التطبيقات ومنصات التوصيل بالمتجر بشكل لحظي وفوري لمعالجة طلبات المزوّدين.",
                "مستهدفة من حجم سوق التوصيل والطلبات في المملكة، مدعومة بنشر القائمة الإلكترونية ومزامنة المخزون وتكامل القنوات.",
                "التوسع بامتلاك حصص أوسع في قنوات الطلبات الإلكترونية مع زيادة عدد المزوّدين النشطين على منصة Thawani."
            ]
            
            colors = [DARK_BLUE, ORANGE, RGBColor(59, 130, 246)] # Azure Blue for final
            
            for j in range(3):
                left_x = positions[j]
                create_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left_x, top_y, box_width, height_y, WHITE, colors[j], 1)
                
                # Colored Indicator Badge at top
                create_shape(slide, MSO_SHAPE.RECTANGLE, left_x + Inches(0.15), top_y + Inches(0.15), box_width - Inches(0.3), Inches(0.15), colors[j])
                
                badge_box = slide.shapes.add_textbox(left_x + Inches(0.15), top_y + Inches(0.4), box_width - Inches(0.3), height_y - Inches(0.55))
                btf = badge_box.text_frame
                btf.word_wrap = True
                
                # Metric
                add_paragraph_formatted(btf, kpi_metrics[j], 18, colors[j], bold=True, align=PP_ALIGN.CENTER)
                btf.add_paragraph()
                
                # Title
                add_paragraph_formatted(btf, kpi_titles[j], 13, DARK_BLUE, bold=True, align=PP_ALIGN.RIGHT)
                btf.add_paragraph()
                
                # Body
                add_paragraph_formatted(btf, kpi_texts[j], 11, MUTED_GREY, align=PP_ALIGN.RIGHT)
                
                # Bullet proof point highlight
                if j == 1 and len(bullets) > 1:
                    btf.add_paragraph()
                    add_paragraph_formatted(btf, "ميزة مسبقة: " + bullets[1][:80] + "...", 10, ORANGE, italic=True)

        # 7. Slide 9 (Wameed Market Share Progress Bars)
        elif slide_no == '9':
            add_header(slide, slide_no, slide_title, slide_sub)
            
            # Left has progress bars. Right has text bullets.
            # Right Column: Bullet text
            bullets_width = Inches(6.1)
            bullets_box = slide.shapes.add_textbox(Inches(6.6), Inches(1.8), bullets_width, Inches(5.0))
            btf = bullets_box.text_frame
            btf.word_wrap = True
            
            add_paragraph_formatted(btf, "أهداف التوسع والمستهدوف التشغيلي", 14, ORANGE, bold=True)
            btf.add_paragraph()
            
            for b in bullets:
                add_paragraph_formatted(btf, "• " + b, 12, DARK_BLUE)
                p_spacer = btf.add_paragraph()
                p_spacer.font.size = Pt(6) # small spacing
                
            # Left Column: Beautiful program progress bars
            left_pos = Inches(0.66)
            top_bars = Inches(1.8)
            width_bars = Inches(5.5)
            height_bars = Inches(5.0)
            
            create_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_bars, width_bars, height_bars, WHITE, DARK_BLUE, 1)
            
            bars_box = slide.shapes.add_textbox(left_pos + Inches(0.2), top_bars + Inches(0.2), width_bars - Inches(0.4), height_bars - Inches(0.4))
            tf_b = bars_box.text_frame
            tf_b.word_wrap = True
            
            add_paragraph_formatted(tf_b, "النسب المئوية المستهدفة للحصة السوقية", 14, DARK_BLUE, bold=True, align=PP_ALIGN.CENTER)
            tf_b.add_paragraph()
            
            # Segment 1: Retail & Supermarket POS
            add_paragraph_formatted(tf_b, "قطاع التجزئة والسوبر ماركت (المستهدف: 40-45%)", 11, ORANGE, bold=True)
            tf_b.add_paragraph() # space-holder for progress bar graphic
            pbar_y1 = top_bars + Inches(1.2)
            # Gray bar (background)
            create_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left_pos + Inches(0.2), pbar_y1, Inches(4.5), Inches(0.25), RGBColor(220, 220, 220))
            # Orange filled bar (45%)
            create_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left_pos + Inches(0.2), pbar_y1, Inches(4.5 * 0.45), Inches(0.25), ORANGE)
            
            tf_b.add_paragraph() # Spacing
            tf_b.add_paragraph()
            
            # Segment 2: Restaurants & Cafes POS
            add_paragraph_formatted(tf_b, "قطاع الضيافة والمطاعم (المستهدف: ~10%)", 11, DARK_BLUE, bold=True)
            tf_b.add_paragraph()
            pbar_y2 = pbar_y1 + Inches(1.4)
            # Gray bar
            create_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left_pos + Inches(0.2), pbar_y2, Inches(4.5), Inches(0.25), RGBColor(220, 220, 220))
            # Yellow filled bar (10%)
            create_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left_pos + Inches(0.2), pbar_y2, Inches(4.5 * 0.12), Inches(0.25), YELLOW)
            
            tf_b.add_paragraph() # Spacing
            tf_b.add_paragraph()
            
            # Segment 3: Specialized Sectors (Bakery, Pharmacy, etc.)
            add_paragraph_formatted(tf_b, "القطاعات والصناعات المتخصصة الأخرى (دخول نشط)", 11, MUTED_GREY, bold=True)
            tf_b.add_paragraph()
            pbar_y3 = pbar_y2 + Inches(1.4)
            # Gray bar
            create_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left_pos + Inches(0.2), pbar_y3, Inches(4.5), Inches(0.25), RGBColor(220, 220, 220))
            # Blue filled bar (25%)
            create_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left_pos + Inches(0.2), pbar_y3, Inches(4.5 * 0.25), Inches(0.25), RGBColor(59, 130, 246))

        # 8. Standard Content Slides (Slides 5, 6, 8, 10 to 19)
        else:
            add_header(slide, slide_no, slide_title, slide_sub)
            
            # RTL layout: Right is Text Bullets, Left is a beautifully styled Screen Visual Blueprint Box.
            text_col_left = Inches(6.5)
            text_col_width = Inches(6.16)
            top_y = Inches(1.8)
            height_y = Inches(5.0)
            
            # Draw Bullets on the Right column
            right_textbox = slide.shapes.add_textbox(text_col_left, top_y, text_col_width, height_y)
            rtf = right_textbox.text_frame
            rtf.word_wrap = True
            
            add_paragraph_formatted(rtf, "الميزات والعمليات الرئيسية", 14, ORANGE, bold=True)
            rtf.add_paragraph()
            
            for b in bullets:
                # Beautiful bullet symbols in Wameed colors
                add_paragraph_formatted(rtf, "• " + b, 12, DARK_BLUE)
                # tiny spacing paragraph
                p_space = rtf.add_paragraph()
                p_space.font.size = Pt(6)
                
            # Draw Visual Screen blueprint frame on the Left column
            visual_col_left = Inches(0.66)
            visual_col_width = Inches(5.5)
            
            # Container border
            create_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, visual_col_left, top_y, visual_col_width, height_y, WHITE, ORANGE if "Thawani" in slide_title or int(slide_no) in [5, 6] else DARK_BLUE, 1)
            
            # Subtle top bar representing app screen header
            header_color = ORANGE if "Thawani" in slide_title or int(slide_no) in [5, 6] else DARK_BLUE
            create_shape(slide, MSO_SHAPE.RECTANGLE, visual_col_left + Inches(0.15), top_y + Inches(0.15), visual_col_width - Inches(0.3), Inches(0.4), header_color)
            
            header_txt_box = slide.shapes.add_textbox(visual_col_left + Inches(0.15), top_y + Inches(0.12), visual_col_width - Inches(0.3), Inches(0.4))
            htf = header_txt_box.text_frame
            add_paragraph_formatted(htf, "محاكاة واجهة التطبيق | App UI Mockup", 10, WHITE, bold=True, align=PP_ALIGN.CENTER)
            
            # Central Inner dashboard frame
            create_shape(slide, MSO_SHAPE.RECTANGLE, visual_col_left + Inches(0.25), top_y + Inches(0.7), visual_col_width - Inches(0.5), Inches(4.0), BG_LIGHT, MUTED_GREY, 1)
            
            inner_textbox = slide.shapes.add_textbox(visual_col_left + Inches(0.35), top_y + Inches(0.85), visual_col_width - Inches(0.7), Inches(3.7))
            itf = inner_textbox.text_frame
            itf.word_wrap = True
            
            # Text inside visual frame to assist builder
            add_paragraph_formatted(itf, "[لقطة شاشة نشطة من الكود - Mockup Stage Screen]", 11, ORANGE if header_color == ORANGE else DARK_BLUE, bold=True, align=PP_ALIGN.CENTER)
            itf.add_paragraph()
            
            add_paragraph_formatted(itf, "مخطط التموضع والتوجيه البصري المقترح:", 11, DARK_BLUE, bold=True, align=PP_ALIGN.RIGHT)
            itf.add_paragraph()
            
            visual_text = sd.get('visual', 'واجهة تشغيلية حقيقية توضح تفاصيل الوحدة والخطوات التشغيلية.')
            add_paragraph_formatted(itf, visual_text, 10, MUTED_GREY, italic=True, align=PP_ALIGN.RIGHT)
            itf.add_paragraph()
            
            # Quick directory guide tag in UI mockups
            if sd.get('guide'):
                itf.add_paragraph()
                add_paragraph_formatted(itf, f"ملف التحقق بالكود:\n{sd['guide']}", 8, DARK_BLUE, italic=True, align=PP_ALIGN.RIGHT)
                
    # Save Presentation
    print(f"Saving compiled PPTX presentation to: {output_path}")
    prs.save(output_path)
    print("Presentation compiled successfully!")

if __name__ == "__main__":
    main()
